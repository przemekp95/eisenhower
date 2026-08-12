from pathlib import Path

from docx import Document
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter
from pypdf.generic import DictionaryObject, NameObject, DecodedStreamObject


def generate_pdf(path: Path) -> None:
  writer = PdfWriter()
  page = writer.add_blank_page(width=595, height=842)
  font = DictionaryObject({
    NameObject("/Type"): NameObject("/Font"),
    NameObject("/Subtype"): NameObject("/Type1"),
    NameObject("/BaseFont"): NameObject("/Helvetica"),
  })
  font_ref = writer._add_object(font)  # pylint: disable=protected-access
  page[NameObject("/Resources")] = DictionaryObject({
    NameObject("/Font"): DictionaryObject({NameObject("/F1"): font_ref}),
  })
  content = DecodedStreamObject()
  content.set_data(
    b"BT /F1 18 Tf 72 760 Td (Corpus PDF validation) Tj "
    b"0 -32 Td /F1 11 Tf (Approved local source. Human OCR gate. Fail closed.) Tj ET"
  )
  page[NameObject("/Contents")] = writer._add_object(content)  # pylint: disable=protected-access
  with path.open("wb") as output:
    writer.write(output)


def generate_ocr_pdf(path: Path) -> None:
  image = Image.new("RGB", (1800, 600), "white")
  draw = ImageDraw.Draw(image)
  font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 96)
  draw.multiline_text(
    (100, 120),
    "OCR VALIDATION\nHUMAN APPROVAL REQUIRED",
    fill="black",
    font=font,
    spacing=50,
  )
  image.save(path, "PDF", resolution=300.0)


def generate_docx(path: Path) -> None:
  document = Document()
  document.add_heading("Walidacja dokumentu DOCX", level=1)
  document.add_paragraph("Zatwierdzone lokalne źródło bez danych użytkownika.")
  document.add_paragraph("Sprawdź pochodzenie", style="List Bullet")
  table = document.add_table(rows=2, cols=2)
  table.cell(0, 0).text = "Bramka"
  table.cell(0, 1).text = "Wynik"
  table.cell(1, 0).text = "OCR"
  table.cell(1, 1).text = "Kontrola człowieka"
  document.save(path)


def generate_pptx(path: Path) -> None:
  presentation = Presentation()
  slide = presentation.slides.add_slide(presentation.slide_layouts[5])
  slide.shapes.title.text = "PPTX extraction validation"
  text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))
  text_box.text_frame.text = "Approved local source. Preserve slide provenance and fail closed."
  table = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(6), Inches(1.5)).table
  table.cell(0, 0).text = "Gate"
  table.cell(0, 1).text = "Status"
  table.cell(1, 0).text = "Source"
  table.cell(1, 1).text = "Approved"
  presentation.save(path)


def main() -> None:
  root = Path(__file__).resolve().parents[2] / "corpus" / "approved-documents"
  root.mkdir(parents=True, exist_ok=True)
  generate_pdf(root / "extraction-golden-pdf.pdf")
  generate_ocr_pdf(root / "extraction-golden-ocr.pdf")
  generate_docx(root / "extraction-golden-docx.docx")
  generate_pptx(root / "extraction-golden-pptx.pptx")


if __name__ == "__main__":
  main()
